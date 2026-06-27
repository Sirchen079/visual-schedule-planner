from __future__ import annotations

import base64
import json
import mimetypes
import re
import secrets
import time
from dataclasses import asdict, dataclass
from pathlib import Path
from typing import Iterable

from fastapi import HTTPException, UploadFile, status
from sqlalchemy.orm import Session

from app.config import settings
from app.models import File
from app.services import file_service

_CHUNK = 1024 * 1024

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".gif"}
TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".markdown",
    ".csv",
    ".json",
    ".jsonl",
    ".html",
    ".htm",
    ".xml",
    ".yaml",
    ".yml",
    ".log",
    ".py",
    ".js",
    ".ts",
    ".tsx",
    ".jsx",
    ".vue",
    ".css",
    ".scss",
    ".less",
    ".sql",
    ".ini",
    ".toml",
    ".tex",
}
PDF_EXTENSIONS = {".pdf"}
OFFICE_EXTENSIONS = {".docx", ".xlsx", ".xlsm", ".pptx"}


@dataclass
class ChatAttachment:
    id: str
    original_name: str
    storage_path: str
    size: int
    mime_type: str
    kind: str
    created_at: float


class TextBudget:
    def __init__(self, limit: int):
        self.limit = limit
        self.used = 0
        self.parts: list[str] = []
        self.truncated = False

    def add(self, text: str) -> None:
        if not text or self.truncated:
            return
        remain = self.limit - self.used
        if remain <= 0:
            self.truncated = True
            return
        if len(text) > remain:
            self.parts.append(text[:remain])
            self.used += remain
            self.truncated = True
            return
        self.parts.append(text)
        self.used += len(text)

    def value(self) -> str:
        text = "".join(self.parts).strip()
        if self.truncated:
            text = f"{text}\n\n[内容过长，已按 {self.limit} 字符上限截断]"
        return text


def save_upload(upload: UploadFile) -> ChatAttachment:
    settings.ai_attachments_dir.mkdir(parents=True, exist_ok=True)
    safe_name = _safe_name(upload.filename or "unnamed")
    attachment_id = secrets.token_urlsafe(24)
    stored_name = f"{attachment_id}-{safe_name}"
    destination = settings.ai_attachments_dir / stored_name
    limit = settings.max_ai_attachment_bytes

    written = 0
    with destination.open("wb") as out:
        while chunk := upload.file.read(_CHUNK):
            written += len(chunk)
            if written > limit:
                out.close()
                destination.unlink(missing_ok=True)
                raise HTTPException(
                    status_code=status.HTTP_413_CONTENT_TOO_LARGE,
                    detail=f"AI 对话附件超过上限 {settings.max_ai_attachment_mb}MB",
                )
            out.write(chunk)

    mime_type = _effective_mime_type(safe_name, upload.content_type)
    attachment = ChatAttachment(
        id=attachment_id,
        original_name=safe_name,
        storage_path=_relative_storage_path(destination),
        size=destination.stat().st_size,
        mime_type=mime_type,
        kind=classify_file(safe_name, mime_type),
        created_at=time.time(),
    )
    _metadata_path(attachment_id).write_text(
        json.dumps(asdict(attachment), ensure_ascii=False),
        encoding="utf-8",
    )
    return attachment


def get_attachment(attachment_id: str) -> ChatAttachment | None:
    if not _valid_attachment_id(attachment_id):
        return None
    path = _metadata_path(attachment_id)
    if not path.exists():
        return None
    try:
        raw = json.loads(path.read_text(encoding="utf-8"))
        attachment = ChatAttachment(**raw)
    except (OSError, TypeError, json.JSONDecodeError):
        return None
    if not content_path(attachment).exists():
        return None
    return attachment


def build_model_attachments(attachment_ids: Iterable[str]) -> list[dict]:
    items = []
    for attachment_id in attachment_ids:
        attachment = get_attachment(attachment_id)
        if attachment is None:
            raise HTTPException(status_code=404, detail=f"对话附件不存在: {attachment_id}")
        path = content_path(attachment)
        if attachment.kind == "image":
            if attachment.size > settings.max_ai_inline_image_bytes:
                raise HTTPException(
                    status_code=status.HTTP_413_CONTENT_TOO_LARGE,
                    detail=f"图片附件超过识图上限 {settings.max_ai_inline_image_mb}MB: {attachment.original_name}",
                )
            items.append(
                {
                    "id": attachment.id,
                    "kind": "image",
                    "filename": attachment.original_name,
                    "mime_type": attachment.mime_type,
                    "size": attachment.size,
                    "data": base64.b64encode(path.read_bytes()).decode("ascii"),
                }
            )
        else:
            try:
                text = extract_document_text(
                    path, attachment.original_name, attachment.mime_type
                )
            except Exception as exc:
                text = (
                    f"附件 {attachment.original_name} 解析失败：{exc.__class__.__name__}。"
                    "请确认文件未损坏，或转换为 PDF、docx、xlsx、pptx、txt 后重试。"
                )
            items.append(
                {
                    "id": attachment.id,
                    "kind": "document",
                    "filename": attachment.original_name,
                    "mime_type": attachment.mime_type,
                    "size": attachment.size,
                    "text": text,
                }
            )
    return items


def save_to_library(
    db: Session,
    attachment_id: str,
    notes: str = "",
    task_id: int | None = None,
) -> File:
    attachment = get_attachment(attachment_id)
    if attachment is None:
        raise ValueError(f"对话附件不存在: {attachment_id}")
    db_file = file_service.save_local_file(
        db,
        content_path(attachment),
        attachment.original_name,
        attachment.mime_type,
        notes or "由 AI 从对话附件保存到资料库",
    )
    if task_id is not None:
        file_service.attach_to_task(db, task_id, db_file.id)
    return db_file


def extract_document_text(path: Path, filename: str, mime_type: str) -> str:
    ext = Path(filename).suffix.lower()
    if ext in PDF_EXTENSIONS or mime_type == "application/pdf":
        return _extract_pdf(path)
    if ext == ".docx":
        return _extract_docx(path)
    if ext in {".xlsx", ".xlsm"}:
        return _extract_xlsx(path)
    if ext == ".pptx":
        return _extract_pptx(path)
    if _is_text_file(ext, mime_type):
        return _extract_plain_text(path)
    return (
        f"附件 {filename} 已上传，但当前不支持解析该格式的正文。"
        "如需长期保存，可让助手把它加入资料库；如需分析内容，请转换为 PDF、docx、xlsx、pptx 或文本格式。"
    )


def classify_file(filename: str, mime_type: str) -> str:
    ext = Path(filename).suffix.lower()
    if ext in IMAGE_EXTENSIONS or mime_type.startswith("image/"):
        return "image"
    if ext in TEXT_EXTENSIONS | PDF_EXTENSIONS | OFFICE_EXTENSIONS or _is_text_file(ext, mime_type):
        return "document"
    return "document"


def content_path(attachment: ChatAttachment) -> Path:
    return settings.database_dir.parent / attachment.storage_path


def _extract_plain_text(path: Path) -> str:
    raw = path.read_bytes()
    for encoding in ("utf-8-sig", "utf-16", "gb18030"):
        try:
            return _truncate(raw.decode(encoding))
        except UnicodeDecodeError:
            continue
    return _truncate(raw.decode("utf-8", errors="replace"))


def _extract_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    if getattr(reader, "is_encrypted", False):
        try:
            reader.decrypt("")
        except Exception:
            return "PDF 已加密，无法读取正文。"
    budget = TextBudget(settings.max_ai_text_chars)
    for index, page in enumerate(reader.pages, start=1):
        budget.add(f"\n\n[PDF 第 {index} 页]\n")
        try:
            budget.add(page.extract_text() or "")
        except Exception as exc:
            budget.add(f"\n[第 {index} 页解析失败: {exc.__class__.__name__}]\n")
        if budget.truncated:
            break
    return budget.value() or "PDF 未提取到可读文本，可能是扫描版或图片型论文。"


def _extract_docx(path: Path) -> str:
    from docx import Document

    document = Document(str(path))
    budget = TextBudget(settings.max_ai_text_chars)
    for paragraph in document.paragraphs:
        budget.add(paragraph.text)
        budget.add("\n")
        if budget.truncated:
            break
    for table_index, table in enumerate(document.tables, start=1):
        if budget.truncated:
            break
        budget.add(f"\n[Word 表格 {table_index}]\n")
        for row in table.rows:
            budget.add("\t".join(cell.text.strip() for cell in row.cells))
            budget.add("\n")
            if budget.truncated:
                break
    return budget.value() or "Word 文档未提取到可读文本。"


def _extract_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(str(path), read_only=True, data_only=True)
    budget = TextBudget(settings.max_ai_text_chars)
    try:
        for sheet in workbook.worksheets:
            budget.add(f"\n[Excel 工作表: {sheet.title}]\n")
            for row in sheet.iter_rows(values_only=True):
                values = ["" if value is None else str(value) for value in row]
                if any(value.strip() for value in values):
                    budget.add("\t".join(values))
                    budget.add("\n")
                if budget.truncated:
                    break
            if budget.truncated:
                break
    finally:
        workbook.close()
    return budget.value() or "Excel 文件未提取到可读文本。"


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation

    presentation = Presentation(str(path))
    budget = TextBudget(settings.max_ai_text_chars)
    for index, slide in enumerate(presentation.slides, start=1):
        budget.add(f"\n[PPT 第 {index} 页]\n")
        for shape in slide.shapes:
            text = ""
            if getattr(shape, "has_text_frame", False):
                text = shape.text
            elif getattr(shape, "has_table", False):
                lines = []
                for row in shape.table.rows:
                    lines.append("\t".join(cell.text.strip() for cell in row.cells))
                text = "\n".join(lines)
            if text:
                budget.add(text)
                budget.add("\n")
            if budget.truncated:
                break
        if budget.truncated:
            break
    return budget.value() or "PPT 文件未提取到可读文本。"


def _truncate(text: str) -> str:
    text = _clean_text(text)
    if len(text) <= settings.max_ai_text_chars:
        return text
    return f"{text[: settings.max_ai_text_chars]}\n\n[内容过长，已按 {settings.max_ai_text_chars} 字符上限截断]"


def _clean_text(text: str) -> str:
    return re.sub(r"\n{4,}", "\n\n\n", text.replace("\r\n", "\n").replace("\r", "\n")).strip()


def _is_text_file(ext: str, mime_type: str) -> bool:
    return (
        mime_type.startswith("text/")
        or mime_type
        in {
            "application/json",
            "application/xml",
            "application/yaml",
            "application/x-yaml",
            "application/javascript",
            "application/typescript",
            "application/sql",
            "text/csv",
        }
        or ext in TEXT_EXTENSIONS
    )


def _effective_mime_type(filename: str, content_type: str | None) -> str:
    if content_type and content_type != "application/octet-stream":
        return content_type
    guessed, _ = mimetypes.guess_type(filename)
    return guessed or "application/octet-stream"


def _safe_name(filename: str) -> str:
    name = Path(filename or "unnamed").name
    for ch in '<>:"/\\|?*':
        name = name.replace(ch, "_")
    return name.strip() or "unnamed"


def _metadata_path(attachment_id: str) -> Path:
    return settings.ai_attachments_dir / f"{attachment_id}.json"


def _valid_attachment_id(attachment_id: str) -> bool:
    return bool(re.fullmatch(r"[A-Za-z0-9_-]{8,120}", attachment_id or ""))


def _relative_storage_path(path: Path) -> str:
    try:
        return str(path.relative_to(settings.database_dir.parent))
    except ValueError:
        return str(path)
